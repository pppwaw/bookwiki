from __future__ import annotations

import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from bookwiki.convert.common import clean_markdown, source_id_from_stem


def convert_pptx_to_md(path: str | Path, *, source_id: str | None = None) -> str:
    pptx_path = Path(path)
    resolved_source_id = source_id_from_stem(source_id or pptx_path.stem)
    slides = _read_slides_with_python_pptx(pptx_path) or _read_slides_from_xml(pptx_path)

    blocks = [f"# {pptx_path.stem}"]
    for index, texts in enumerate(slides, start=1):
        title = texts[0] if texts else f"Slide {index}"
        body_lines = texts[1:] if len(texts) > 1 else []
        source_ref = f"{resolved_source_id}-slide{index:02d}"
        body = "\n".join(f"- {line}" for line in body_lines)
        blocks.append(
            clean_markdown(
                f"## Slide {index}: {title}\n\n<!-- source_ref: {source_ref} -->\n\n{body}"
            )
        )
    return "\n\n".join(blocks).strip() + "\n"


def _read_slides_with_python_pptx(path: Path) -> list[list[str]]:
    try:
        from pptx import Presentation
    except ImportError:
        return []

    presentation = Presentation(str(path))
    slides: list[list[str]] = []
    for slide in presentation.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = clean_markdown(str(shape.text))
            if text:
                texts.extend(line.strip() for line in text.splitlines() if line.strip())
        slides.append(texts)
    return slides


def _read_slides_from_xml(path: Path) -> list[list[str]]:
    slides: list[list[str]] = []
    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            (
                name
                for name in archive.namelist()
                if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
            ),
            key=_slide_sort_key,
        )
        for slide_name in slide_names:
            root = ElementTree.fromstring(archive.read(slide_name))
            texts = [
                node.text.strip()
                for node in root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                if node.text and node.text.strip()
            ]
            slides.append(texts)
    return slides


def _slide_sort_key(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0
